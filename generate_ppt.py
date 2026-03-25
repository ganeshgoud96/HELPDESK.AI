from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Optional: We could set slide width/height for 16:9, but default 4:3 is standard unless specified.
    # Let's use 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Layouts
    TITLE_SLIDE = 0
    BULLET_SLIDE = 1
    TWO_CONTENT = 3
    TITLE_ONLY = 5

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "HELPDESK.AI\nAI-Powered Ticket Creation & Categorization"
    subtitle.text = "Infosys Springboard Virtual Internship 6.0\nTeam Presentation\nMarch 2026"

    # 2. Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Problem Statement"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Manual Ticket Triage Bottlenecks:"
    p = tf.add_paragraph()
    p.text = "Manual Creation: Takes 4-8 minutes per ticket."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Inconsistent Routing: Human agents often mistriage severity and teams."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Time Consuming: Slower response times breach SLAs."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Ticket Floods: Outages cause hundreds of redundant manual tickets."
    p.level = 1

    # 3. Project Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Project Objectives"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Core Automation Goals:"
    p = tf.add_paragraph()
    p.text = "Automatic Ticket Creation & Categorization via AI."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "NER (Named Entity Recognition) for automated metadata extraction."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Semantic Duplicate Detection using Cosine Similarity."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Real-time Chat Troubleshooting via Google Gemini."
    p.level = 1

    # 4. System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    slide.shapes.title.text = "System Architecture & Flow"
    # Placeholder for diagram
    txbox = slide.shapes.add_textbox(Inches(3), Inches(2.5), Inches(7.33), Inches(2))
    tf = txbox.text_frame
    tf.text = "[ 🖼️ INSERT ARCHITECTURE / FLOW DIAGRAM HERE ]\n\nFastAPI ➔ DistilBERT / NER / Gemini ➔ Supabase ➔ React/Vite"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for par in tf.paragraphs:
        par.font.size = Pt(24)
        par.font.color.rgb = RGBColor(100, 100, 100)

    # 5. Dataset Overview
    slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT])
    slide.shapes.title.text = "Dataset Overview"
    left_frame = slide.placeholders[1].text_frame
    right_frame = slide.placeholders[2].text_frame

    left_frame.text = "Final_Balanced_10000_IT_Support_Tickets.csv"
    p = left_frame.add_paragraph()
    p.text = "10,000 synthetic IT support tickets."
    p.level = 1
    p = left_frame.add_paragraph()
    p.text = "5 Predicted Fields: Category, Subcategory, Priority, Auto-Resolve, Assigned Team."
    p.level = 1
    p = left_frame.add_paragraph()
    p.text = "Challenges: Maintaining class balance across rare hardware vs standard software issues."
    p.level = 1

    right_frame.text = "[ 📊 INSERT PIE CHART HERE ]\n\n(Data Distribution / Class Balance Chart)"
    for par in right_frame.paragraphs:
        par.font.size = Pt(24)
        par.font.color.rgb = RGBColor(100, 100, 100)
    right_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 6. Data Preprocessing
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Data Preprocessing"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Pipeline Steps:"
    p = tf.add_paragraph()
    p.text = "Cleaning: Strict null dropping & lowercase standardization."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Hinglish Handling: Mixed language normalized via specific dataset expansion."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Tokenization & Lemmatization: HuggingFace BertTokenizerFast (Max Length 256)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Encoding: Scikit-Learn LabelEncoder for all 5 categorical targets."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "[ 🖼️ INSERT BEFORE/AFTER DATAFRAME SCREENSHOT ]"
    p.level = 0
    p.font.color.rgb = RGBColor(0, 100, 200)

    # 7. Model Development
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Model Development: DistilBERT v3"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Model Journey:"
    p = tf.add_paragraph()
    p.text = "Models Tried: Initial Scikit-Learn pipelines, basic LSTM, fine-tuned DistilBERT v2."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Model Finalized: Custom MultiOutputClassifier (bert-base-uncased)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Why?: Joint representation learning across 5 heads prevented conflicting predictions and handled deep context."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Evaluation Scores: ~94% overall accuracy across the 5 output metrics."
    p.level = 1

    # 8. NER (Named Entity Recognition)
    slide = prs.slides.add_slide(prs.slide_layouts[TWO_CONTENT])
    slide.shapes.title.text = "Named Entity Recognition (NER)"
    left_frame = slide.placeholders[1].text_frame
    right_frame = slide.placeholders[2].text_frame

    left_frame.text = "Automated Extraction:"
    p = left_frame.add_paragraph()
    p.text = "Approach: Transformer-based entity harvesting."
    p.level = 1
    p = left_frame.add_paragraph()
    p.text = "Extracts: Device names, Hostnames, IP addresses, OS versions, specific Error codes."
    p.level = 1

    right_frame.text = "JSON Output Example:\n{\n  'Device': 'MacBook Pro M2',\n  'IP': '192.168.1.14',\n  'OS': 'macOS 14.2'\n}"
    for par in right_frame.paragraphs:
        par.font.size = Pt(18)
        par.font.name = 'Courier New'

    # 9. Ticket Generation Engine
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Ticket Generation Engine"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "How Everything is Combined:"
    p = tf.add_paragraph()
    p.text = "Step 1: User submits text/screenshot."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Step 2: Optical Character Recognition (OCR) reads the image."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Step 3: DistilBERT V3 predicts Category, Priority, and Team."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Step 4: NER extracts underlying metadata."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Step 5: Sentence-Transformers check for Duplicates before Supabase insertion."
    p.level = 1

    # 10. Integration & Deployment
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Integration & Deployment"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "End-to-End Orchestration:"
    p = tf.add_paragraph()
    p.text = "Frontend: React 19, Vite, Tailwind CSS (Deployed on Vercel)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend: Python 3.12, PyTorch, FastAPI (Deployed on Hugging Face Spaces)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Database: PostgreSQL via Supabase with strict Row-Level Security."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "API Workflow: Strictly decoupled JSON REST APIs communicating between Vercel Edge and HuggingFace CPUs."
    p.level = 1

    # 11. Demo
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    slide.shapes.title.text = "Live Demo"
    txbox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7.33), Inches(2))
    tf = txbox.text_frame
    tf.text = "[ 🖥️ Switch to Live Application/Video Demo Here ]"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = RGBColor(200, 50, 50)

    # 12. Challenges & Learning
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Challenges & Learnings"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Key Obstacles Overcome:"
    p = tf.add_paragraph()
    p.text = "Dataset Class Imbalance: Solved via synthetic injection and stratified train/test splitting."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Model Cold-Starts: Optimized HuggingFace FastAPI lifespan hooks to load models strictly once on startup."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Multi-Lingual context (Hinglish): Overcome by expanding the BERT training set with localized colloquial datasets."
    p.level = 1

    # 13. Conclusion & Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = "Conclusion & Future Scope"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Conclusion:"
    p = tf.add_paragraph()
    p.text = "Successfully engineered a 500ms AI neural pipeline replacing 4-8 minute manual triage limits."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Scope:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Active Learning Loops based on 'Corrections Log' mapping."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Native Mobile Application integration with push-notification SLAs."
    p.level = 1

    # 14. Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    txbox = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.333), Inches(2))
    tf = txbox.text_frame
    tf.text = "Thank You\n\nQuestions & Answers?"
    for par in tf.paragraphs:
        par.alignment = PP_ALIGN.CENTER
        par.font.size = Pt(44)
        par.font.bold = True

    prs.save('docs/Final_Presentation.pptx')
    print("Presentation Generated: docs/Final_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
